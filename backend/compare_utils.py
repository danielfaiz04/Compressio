import os
import mimetypes
import numpy as np
from PIL import Image
from io import BytesIO
import base64
import logging
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from typing import Optional

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def get_file_type(mime_type: str) -> str:
    """
    Determines a simplified file type category from a MIME type.
    """
    if mime_type.startswith('image/'):
        return "image"
    elif mime_type == 'application/pdf':
        return "pdf"
    elif mime_type.startswith('video/'):
        return "video"
    elif mime_type.startswith('audio/'):
        return "audio"
    elif mime_type.startswith('text/') or mime_type == 'application/json':
        return "text"
    elif mime_type in [
        'application/msword',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'application/vnd.ms-excel',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'application/vnd.ms-powerpoint',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    ]:
        return "document"
    elif mime_type in ['application/zip', 'application/x-tar', 'application/x-rar-compressed']:
        return "archive"
    else:
        return "other"

def calculate_compression_ratio(original_size: int, compressed_size: int) -> float:
    """
    Calculates the compression ratio.
    Ratio is 1 - (compressed_size / original_size).
    Returns 0.0 if original_size is zero or if compressed_size is greater than or equal to original_size.
    """
    if original_size == 0:
        return 0.0
    if compressed_size >= original_size:
        return 0.0
    return (1 - (compressed_size / original_size)) * 100

def create_image_preview(file_path: str, file_type: str, size=(128, 128)) -> Optional[str]:
    """
    Generates a base64 encoded PNG thumbnail for various file types.
    """
    if file_type == "image":
        try:
            img = Image.open(file_path).convert("RGB")
            img.thumbnail(size, Image.Resampling.LANCZOS)
            buffered = BytesIO()
            img.save(buffered, format="PNG")
            return base64.b64encode(buffered.getvalue()).decode("utf-8")
        except Exception as e:
            logging.error(f"Error creating image preview for {file_path}: {e}")
            return None
    elif file_type == "pdf":
        try:
            # This requires a PDF rendering library like pdf2image, or a simpler approach
            # For simplicity, we'll just try to get the first page as an image if possible
            # This is a placeholder and might require external tools like poppler
            from pdf2image import convert_from_path
            images = convert_from_path(file_path, first_page=1, last_page=1, size=size)
            if images:
                buffered = BytesIO()
                images[0].save(buffered, format="PNG")
                return base64.b64encode(buffered.getvalue()).decode("utf-8")
            return None
        except ImportError:
            logging.warning("pdf2image not installed. PDF previews will not be generated. Please install it with 'pip install pdf2image'. Also ensure Poppler is installed and in your PATH.")
            return None
        except Exception as e:
            logging.error(f"Error creating PDF preview for {file_path}: {e}")
            return None
    elif file_type == "text":
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(500) # Read first 500 characters
            # Create a simple image from text using Pillow
            from PIL import ImageDraw, ImageFont
            img = Image.new('RGB', (size[0]*2, size[1]*2), color = (255, 255, 255))
            d = ImageDraw.Draw(img)
            try:
                fnt = ImageFont.truetype("arial.ttf", 15) # Adjust font as needed
            except IOError:
                fnt = ImageFont.load_default()
            d.text((10,10), content, fill=(0,0,0), font=fnt)
            buffered = BytesIO()
            img.save(buffered, format="PNG")
            return base64.b64encode(buffered.getvalue()).decode("utf-8")
        except Exception as e:
            logging.error(f"Error creating text preview for {file_path}: {e}")
            return None
    elif file_type == "document":
        try:
            if file_path.endswith('.docx'):
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
            elif file_path.endswith(('.ppt', '.pptx')):
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                text = "\n".join(text)
            else: # For .doc, .xls, .xlsx, assume text extraction is harder without external libs
                return None 

            content = text[:500] # Take first 500 characters
            from PIL import ImageDraw, ImageFont
            img = Image.new('RGB', (size[0]*2, size[1]*2), color = (255, 255, 255))
            d = ImageDraw.Draw(img)
            try:
                fnt = ImageFont.truetype("arial.ttf", 15)
            except IOError:
                fnt = ImageFont.load_default()
            d.text((10,10), content, fill=(0,0,0), font=fnt)
            buffered = BytesIO()
            img.save(buffered, format="PNG")
            return base64.b64encode(buffered.getvalue()).decode("utf-8")
        except Exception as e:
            logging.error(f"Error creating document preview for {file_path}: {e}")
            return None
    
    # For other types, return None or a generic placeholder
    return None 